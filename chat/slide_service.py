"""
Slide Service - Export slides to PPTX, PDF, and other formats
Modern presentation export engine
"""
import os
import io
import base64
from datetime import datetime
from django.conf import settings
from django.core.files.base import ContentFile


def export_to_pptx(deck, slides):
    """
    Export slide deck to PowerPoint (.pptx)
    Uses python-pptx library
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError as e:
        raise ImportError(f"python-pptx import error: {e}. Run: pip install python-pptx")
    
    prs = Presentation()
    
    # Set slide dimensions based on aspect ratio
    if deck.aspect_ratio == '16:9':
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    elif deck.aspect_ratio == '4:3':
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    else:  # 21:9
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(6.857)
    
    # Get theme colors
    theme = deck.get_theme_config()
    
    for slide in slides:
        content = slide.content or {}
        layout = slide.layout
        
        # Create slide
        blank_layout = prs.slide_layouts[6]  # Blank layout
        pptx_slide = prs.slides.add_slide(blank_layout)
        
        # Set background
        bg = pptx_slide.background
        fill = bg.fill
        fill.solid()
        
        # Parse background color
        bg_color = theme['bg']
        if bg_color.startswith('#'):
            bg_rgb = hex_to_rgb(bg_color)
            fill.fore_color.rgb = RGBColor(bg_rgb[0], bg_rgb[1], bg_rgb[2])
        else:
            fill.fore_color.rgb = RGBColor(15, 23, 42)  # Default dark
        
        # Add content based on layout
        if layout == 'title':
            add_title_slide_content(pptx_slide, content, theme)
        elif layout == 'content':
            add_content_slide_content(pptx_slide, content, theme)
        elif layout == 'two-column':
            add_two_column_content(pptx_slide, content, theme)
        elif layout in ['image-text', 'text-image']:
            add_image_text_content(pptx_slide, content, theme, layout)
        elif layout == 'quote':
            add_quote_content(pptx_slide, content, theme)
        elif layout == 'section':
            add_section_content(pptx_slide, content, theme)
        else:
            add_content_slide_content(pptx_slide, content, theme)
    
    # Save to buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    
    return buffer


def add_title_slide_content(slide, content, theme):
    """Add title slide content"""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    
    title = content.get('title', 'Untitled')
    subtitle = content.get('subtitle', '')
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(theme['text'])
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(0.8))
        sub_frame = sub_box.text_frame
        p = sub_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = hex_to_rgb(theme['text'])
        p.alignment = PP_ALIGN.CENTER

    # Add canvas elements
    add_canvas_elements_to_slide(slide, content, theme)
    apply_slide_background_and_effects(slide, content, theme)


def add_content_slide_content(slide, content, theme):
    """Add standard content slide"""
    from pptx.util import Inches, Pt
    
    title = content.get('title', '')
    bullets = content.get('bullets', [])
    body_content = content.get('content', '')
    
    # Title
    if title:
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11), Inches(0.8))
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(theme['text'])
    
    # Content area
    if bullets:
        left = Inches(0.75)
        top = Inches(1.5)
        width = Inches(11)
        height = Inches(5.5)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(24)
            p.font.color.rgb = hex_to_rgb(theme['text'])
            p.space_before = Pt(12)
    elif body_content:
        text_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11), Inches(5.5))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = body_content[:500]  # Limit content
        p.font.size = Pt(20)
        p.font.color.rgb = hex_to_rgb(theme['text'])

    # Add canvas elements
    add_canvas_elements_to_slide(slide, content, theme)
    apply_slide_background_and_effects(slide, content, theme)


def add_two_column_content(slide, content, theme):
    """Add two column layout"""
    from pptx.util import Inches, Pt
    
    title = content.get('title', '')
    left_bullets = content.get('left_bullets', [])
    right_bullets = content.get('right_bullets', [])
    
    if title:
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11), Inches(0.8))
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(theme['text'])
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.25), Inches(5.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for i, bullet in enumerate(left_bullets[:5]):  # Limit to 5
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = hex_to_rgb(theme['text'])
        p.space_before = Pt(10)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.5), Inches(5.25), Inches(5.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    for i, bullet in enumerate(right_bullets[:5]):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = hex_to_rgb(theme['text'])
        p.space_before = Pt(10)

    # Add canvas elements
    add_canvas_elements_to_slide(slide, content, theme)
    apply_slide_background_and_effects(slide, content, theme)


def add_image_text_content(slide, content, theme, layout):
    """Add image + text layout"""
    from pptx.util import Inches, Pt
    
    title = content.get('title', '')
    body_content = content.get('content', '')
    bullets = content.get('bullets', [])
    image = content.get('image', {})
    
    if title:
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11), Inches(0.8))
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(theme['text'])
    
    # Placeholder for image
    image_url = image.get('url', '')
    if layout == 'image-text':
        # Image on left
        if image_url and image_url.startswith(('http://', 'https://')):
            try:
                add_image_from_url(slide, image_url, Inches(0.75), Inches(1.5), Inches(5.5), Inches(5.5))
            except:
                pass
        
        # Text on right
        text_left = Inches(6.5)
    else:
        # Text on left
        text_left = Inches(0.75)
        
        # Image on right
        if image_url and image_url.startswith(('http://', 'https://')):
            try:
                add_image_from_url(slide, image_url, Inches(6.5), Inches(1.5), Inches(5.5), Inches(5.5))
            except:
                pass
    
    # Add text content
    text_box = slide.shapes.add_textbox(text_left, Inches(1.5), Inches(5.25), Inches(5.5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    if body_content:
        p = text_frame.paragraphs[0]
        p.text = body_content[:300]
        p.font.size = Pt(18)
        p.font.color.rgb = hex_to_rgb(theme['text'])
    elif bullets:
        for i, bullet in enumerate(bullets[:4]):
            if i == 0 and not body_content:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(20)
            p.font.color.rgb = hex_to_rgb(theme['text'])
            p.space_before = Pt(8)

    # Add canvas elements
    add_canvas_elements_to_slide(slide, content, theme)
    apply_slide_background_and_effects(slide, content, theme)


def add_quote_content(slide, content, theme):
    """Add quote slide content"""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    
    quote_data = content.get('quote', {})
    quote_text = quote_data.get('text', content.get('content', ''))
    author = quote_data.get('author', content.get('author', ''))
    
    # Quote text
    quote_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(3))
    quote_frame = quote_box.text_frame
    quote_frame.word_wrap = True
    p = quote_frame.paragraphs[0]
    p.text = f'"{quote_text}"'
    p.font.size = Pt(40)
    p.font.italic = True
    p.font.color.rgb = hex_to_rgb(theme['text'])
    p.alignment = PP_ALIGN.CENTER
    
    # Author
    if author:
        author_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11), Inches(0.6))
        p = author_box.text_frame.paragraphs[0]
        p.text = f"— {author}"
        p.font.size = Pt(24)
        p.font.color.rgb = hex_to_rgb(theme['text'])
        p.alignment = PP_ALIGN.CENTER

    # Add canvas elements
    add_canvas_elements_to_slide(slide, content, theme)

    # Apply background and effects
    apply_slide_background_and_effects(slide, content, theme)


def add_section_content(slide, content, theme):
    """Add section divider slide"""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    
    title = content.get('title', '')
    subtitle = content.get('subtitle', '')
    
    if title:
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.2))
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(theme['text'])
        p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(0.8))
        p = sub_box.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = hex_to_rgb(theme['text'])
        p.alignment = PP_ALIGN.CENTER

    # Add canvas elements
    add_canvas_elements_to_slide(slide, content, theme)

    # Apply background and effects
    apply_slide_background_and_effects(slide, content, theme)


def add_image_from_url(slide, url, left, top, width, height):
    """Add image from URL to slide"""
    from pptx.util import Inches
    import requests
    
    try:
        response = requests.get(url, timeout=10, headers={'User-Agent': 'Mozilla/5.0'})
        if response.status_code == 200:
            image_stream = io.BytesIO(response.content)
            slide.shapes.add_picture(image_stream, left, top, width, height)
    except Exception as e:
        print(f"Error adding image: {e}")


def hex_to_rgb(hex_color):
    """Convert hex color to RGBColor object"""
    from pptx.dml.color import RGBColor
    hex_color = hex_color.lstrip('#')
    if len(hex_color) == 3:
        hex_color = ''.join([c*2 for c in hex_color])
    try:
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )
    except:
        return RGBColor(255, 255, 255)  # Default white


def add_canvas_elements_to_slide(slide, content, theme):
    """Add canvas elements (images, shapes, icons) to PPTX slide"""
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    
    # Get canvas elements from content
    canvas_elements = content.get('_canvasElements', []) or content.get('canvasElements', [])
    if not canvas_elements:
        return
    
    # Get slide dimensions based on aspect ratio (default 16:9)
    slide_width_inches = 13.333
    slide_height_inches = 7.5
    
    for element in canvas_elements:
        if not isinstance(element, dict):
            continue
            
        el_type = element.get('type', '')
        position = element.get('position', {'x': 50, 'y': 50})
        
        # Convert percentage position to inches
        left = Inches((position.get('x', 50) / 100) * slide_width_inches)
        top = Inches((position.get('y', 50) / 100) * slide_height_inches)
        
        try:
            if el_type == 'image':
                image_url = element.get('src', '')
                size = element.get('size', {})
                width_pct = size.get('width', '60%')
                if isinstance(width_pct, str) and '%' in width_pct:
                    width_inches = Inches(float(width_pct.replace('%', '')) / 100 * slide_width_inches * 0.6)
                else:
                    width_inches = Inches(3)
                height_inches = Inches(2.5)
                
                if image_url:
                    add_canvas_image_to_slide(slide, image_url, left, top, width_inches, height_inches)
                    
            elif el_type == 'shape':
                shape_type = element.get('shape', 'rectangle')
                fill_color = element.get('fill', '#3b82f6')
                stroke_color = element.get('stroke', 'transparent')
                stroke_width = element.get('strokeWidth', 2)
                add_canvas_shape_to_slide(slide, shape_type, left, top, fill_color, stroke_color, stroke_width)
                
            elif el_type == 'icon':
                icon_type = element.get('icon', 'chart')
                icon_color = element.get('color', '#3b82f6')
                add_canvas_icon_to_slide(slide, icon_type, left, top, icon_color)
        except Exception as e:
            print(f"Error adding canvas element {el_type}: {e}")


def add_canvas_image_to_slide(slide, url, left, top, width, height):
    """Add image from URL to slide at specific position"""
    import requests
    import io
    
    try:
        response = requests.get(url, timeout=10, headers={'User-Agent': 'Mozilla/5.0'})
        if response.status_code == 200:
            image_stream = io.BytesIO(response.content)
            # Center the image around the position point
            actual_left = left - (width / 2)
            actual_top = top - (height / 2)
            if actual_left < 0: actual_left = 0
            if actual_top < 0: actual_top = 0
            slide.shapes.add_picture(image_stream, actual_left, actual_top, width, height)
    except Exception as e:
        print(f"Error adding canvas image from URL: {e}")


def add_canvas_shape_to_slide(slide, shape_type, left, top, fill_color, stroke_color, stroke_width):
    """Add shape to slide"""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    
    shape_width = Inches(1.5)
    shape_height = Inches(1.5)
    
    # Center around position
    actual_left = left - (shape_width / 2)
    actual_top = top - (shape_height / 2)
    if actual_left < 0: actual_left = 0
    if actual_top < 0: actual_top = 0
    
    # Map shape types to MSO_SHAPE
    shape_map = {
        'rectangle': MSO_SHAPE.RECTANGLE,
        'circle': MSO_SHAPE.OVAL,
        'triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,
        'diamond': MSO_SHAPE.DIAMOND,
        'star': MSO_SHAPE.STAR_5_POINT,
        'arrow': MSO_SHAPE.RIGHT_ARROW,
        'heart': MSO_SHAPE.HEART,
        'hexagon': MSO_SHAPE.HEXAGON,
    }
    
    mso_shape = shape_map.get(shape_type, MSO_SHAPE.RECTANGLE)
    
    # Add shape
    shape = slide.shapes.add_shape(mso_shape, actual_left, actual_top, shape_width, shape_height)
    
    # Set fill color
    if fill_color and fill_color != 'transparent':
        try:
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
        except:
            pass
    
    # Set stroke/border
    if stroke_color and stroke_color != 'transparent':
        try:
            shape.line.color.rgb = hex_to_rgb(stroke_color)
            shape.line.width = Pt(stroke_width)
        except:
            pass


def add_canvas_icon_to_slide(slide, icon_type, left, top, color):
    """Add icon as a colored shape placeholder"""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt
    
    icon_size = Inches(0.8)
    actual_left = left - (icon_size / 2)
    actual_top = top - (icon_size / 2)
    if actual_left < 0: actual_left = 0
    if actual_top < 0: actual_top = 0
    
    # Add a colored rectangle as icon placeholder
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, actual_left, actual_top, icon_size, icon_size)
    
    try:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(color)
        shape.line.fill.background()
    except:
        pass
    
    # Add icon label below
    try:
        text_box = slide.shapes.add_textbox(actual_left, actual_top + icon_size + Inches(0.05), icon_size, Inches(0.25))
        p = text_box.text_frame.paragraphs[0]
        p.text = icon_type[:10]
        p.font.size = Pt(7)
        p.font.color.rgb = hex_to_rgb(color)
        p.alignment = 2  # Center alignment
    except:
        pass


def apply_slide_background_and_effects(slide, content, theme):
    """Apply background and effects from canvas elements"""
    # Apply background if specified
    background = content.get('_background') or content.get('background')
    if background and isinstance(background, str):
        if background.startswith('#'):
            try:
                from pptx.dml.color import RGBColor
                bg = slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = hex_to_rgb(background)
            except Exception as e:
                print(f"Error applying background: {e}")


def export_to_pdf(deck, slides, output_path=None):
    """
    Export slide deck to PDF using Playwright
    Requires: pip install playwright
    """
    try:
        from playwright.sync_api import sync_playwright
    except ImportError:
        raise ImportError("playwright not installed. Run: pip install playwright && playwright install chromium")
    
    from .slide_renderer import render_deck_to_html
    
    # Generate HTML
    html_content = render_deck_to_html(deck, slides)
    
    # Create temporary HTML file
    import tempfile
    with tempfile.NamedTemporaryFile(mode='w', suffix='.html', delete=False, encoding='utf-8') as f:
        f.write(html_content)
        html_path = f.name
    
    try:
        # Generate PDF with Playwright
        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page()
            
            # Set viewport based on aspect ratio
            if deck.aspect_ratio == '16:9':
                page.set_viewport_size({'width': 1920, 'height': 1080})
            elif deck.aspect_ratio == '4:3':
                page.set_viewport_size({'width': 1440, 'height': 1080})
            else:  # 21:9
                page.set_viewport_size({'width': 2560, 'height': 1080})
            
            page.goto(f'file://{html_path}')
            page.wait_for_load_state('networkidle')
            
            # Generate PDF
            pdf_buffer = page.pdf(
                format='A4' if deck.aspect_ratio == '4:3' else 'A4',
                landscape=True,
                print_background=True,
                margin={'top': '0', 'right': '0', 'bottom': '0', 'left': '0'}
            )
            
            browser.close()
            
            return io.BytesIO(pdf_buffer)
    finally:
        # Cleanup
        try:
            os.unlink(html_path)
        except:
            pass


def export_to_images(deck, slides, output_dir=None):
    """
    Export slides as image sequence using Playwright
    """
    try:
        from playwright.sync_api import sync_playwright
    except ImportError:
        raise ImportError("playwright not installed")
    
    from .slide_renderer import render_deck_to_html
    
    html_content = render_deck_to_html(deck, slides)
    
    import tempfile
    with tempfile.NamedTemporaryFile(mode='w', suffix='.html', delete=False, encoding='utf-8') as f:
        f.write(html_content)
        html_path = f.name
    
    images = []
    try:
        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page()
            
            if deck.aspect_ratio == '16:9':
                page.set_viewport_size({'width': 1920, 'height': 1080})
            elif deck.aspect_ratio == '4:3':
                page.set_viewport_size({'width': 1440, 'height': 1080})
            else:
                page.set_viewport_size({'width': 2560, 'height': 1080})
            
            page.goto(f'file://{html_path}')
            page.wait_for_load_state('networkidle')
            
            # Take screenshot of each slide
            slide_elements = page.query_selector_all('.slide-page')
            for i, slide_elem in enumerate(slide_elements):
                screenshot = slide_elem.screenshot(type='png')
                images.append({
                    'index': i + 1,
                    'data': screenshot,
                    'filename': f'slide_{i+1:03d}.png'
                })
            
            browser.close()
    finally:
        try:
            os.unlink(html_path)
        except:
            pass
    
    return images


class SlideService:
    """
    Main service class for slide operations
    """
    
    @staticmethod
    def create_deck(user, title, description='', theme='modern', aspect_ratio='16:9'):
        """Create a new slide deck"""
        from .models import SlideDeck
        return SlideDeck.objects.create(
            user=user,
            title=title,
            description=description,
            theme=theme,
            aspect_ratio=aspect_ratio
        )
    
    @staticmethod
    def add_slide(deck, layout='content', order=None, content=None):
        """Add a slide to a deck"""
        from .models import Slide
        if order is None:
            order = deck.slides.count()
        return Slide.objects.create(
            deck=deck,
            layout=layout,
            order=order,
            content=content or {}
        )
    
    @staticmethod
    def reorder_slides(deck, slide_orders):
        """Reorder slides: slide_orders = [(slide_id, new_order), ...]"""
        from .models import Slide
        for slide_id, new_order in slide_orders:
            Slide.objects.filter(id=slide_id, deck=deck).update(order=new_order)
    
    @staticmethod
    def duplicate_slide(slide):
        """Duplicate a slide"""
        from .models import Slide
        new_order = slide.order + 1
        # Shift other slides
        Slide.objects.filter(deck=slide.deck, order__gte=new_order).update(order=models.F('order') + 1)
        
        return Slide.objects.create(
            deck=slide.deck,
            layout=slide.layout,
            order=new_order,
            content=slide.content.copy() if slide.content else {}
        )
    
    @staticmethod
    def export_deck(deck, format_type='pptx'):
        """
        Export deck to specified format
        Returns: (filename, buffer)
        """
        slides = deck.slides.all()
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        
        if format_type == 'pptx':
            buffer = export_to_pptx(deck, slides)
            filename = f"{deck.title.replace(' ', '_')}_{timestamp}.pptx"
            return filename, buffer
        
        elif format_type == 'pdf':
            buffer = export_to_pdf(deck, slides)
            filename = f"{deck.title.replace(' ', '_')}_{timestamp}.pdf"
            return filename, buffer
        
        elif format_type == 'html':
            from .slide_renderer import render_deck_to_html
            html = render_deck_to_html(deck, slides)
            buffer = io.BytesIO(html.encode('utf-8'))
            filename = f"{deck.title.replace(' ', '_')}_{timestamp}.html"
            return filename, buffer
        
        elif format_type == 'images':
            images = export_to_images(deck, slides)
            # Return as ZIP
            import zipfile
            buffer = io.BytesIO()
            with zipfile.ZipFile(buffer, 'w', zipfile.ZIP_DEFLATED) as zf:
                for img in images:
                    zf.writestr(img['filename'], img['data'])
            buffer.seek(0)
            filename = f"{deck.title.replace(' ', '_')}_{timestamp}_images.zip"
            return filename, buffer
        
        else:
            raise ValueError(f"Unsupported format: {format_type}")


# Import models at end to avoid circular import
from django.db import models
